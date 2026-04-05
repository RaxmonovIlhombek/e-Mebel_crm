import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
image_slide_layout = prs.slide_layouts[5]

# Define image paths
logo_2 = r"C:\Users\rahmo\OneDrive\Desktop\e-mebel\emebel-crm-v2\public\logo (2).png"
logo_icon = r"C:\Users\rahmo\OneDrive\Desktop\e-mebel\emebel-crm-v2\public\logo_icon.png"
logo_yozuv = r"C:\Users\rahmo\OneDrive\Desktop\e-mebel\emebel-crm-v2\public\logo_yozuv.png"

def add_logo(slide, logo_path, left, top, height):
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, left, top, height=height)

# --- Slide 1: Title Slide ---
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "E-Mebel CRM"
subtitle.text = "Mebel ishlab chiqarishni avtomatlashtirish va mijozlar bilan ishlash tizimi\n\nTaqdimot"

# Add logo to title slide
if os.path.exists(logo_2):
    slide.shapes.add_picture(logo_2, left=Inches(3.5), top=Inches(0.5), height=Inches(2))

# --- Slide 2: Loyiha maqsadi va vazifalari ---
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Loyiha maqsadi va afzalliklari"
tf = body_shape.text_frame
tf.text = "E-Mebel CRM tizimi quyidagi asosiy muammolarni hal qiladi:"

p = tf.add_paragraph()
p.text = "Mijozlar buyurtmalarini boshqarishni to'liq raqamlashtirish"
p.level = 1

p = tf.add_paragraph()
p.text = "Omborxona va xomashyo hisobini aniq yuritish"
p.level = 1

p = tf.add_paragraph()
p.text = "Moliya va to'lovlarni nazorat qilish (qarzlar, to'lovlar, xarajatlar)"
p.level = 1

p = tf.add_paragraph()
p.text = "Xodimlarning ruxsatlarini boshqarish (Admin, Menejer, Omborchi, Kassir)"
p.level = 1

add_logo(slide, logo_yozuv, Inches(7.5), Inches(0.2), Inches(0.7))

# --- Slide 3: Tizim imkoniyatlari ---
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Asosiy bo'limlar va Modullar"
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = "Analitika (Dashboard): Barcha muhim ko'rsatkichlar bitta ekranda"
p.level = 0
p = tf.add_paragraph()
p.text = "Buyurtmalar: Yangi buyurtma qo'shish, statusni o'zgartirish va PDF shartnoma yuklab olish"
p.level = 0
p = tf.add_paragraph()
p.text = "Mijozlar Portali: Mijozlar o'z akkauntlariga kirib qarzlarini va buyurtma holatini kuzatishi"
p.level = 0
p = tf.add_paragraph()
p.text = "Ombor: Xomashyo va tayyor mahsulotlarni nazorat qilish"
p.level = 0

add_logo(slide, logo_icon, Inches(8.5), Inches(0.2), Inches(0.8))

# --- Slide 4: Telegram Bot Integratsiyasi ---
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Telegram Bot va Avtomatlashtirish"
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = "Telegram Bot orqali tizimga kirish:"
p.level = 0
p = tf.add_paragraph()
p.text = "Xodimlar o'z rollariga mos keluvchi tugmalar va bildirishnomalar oladi"
p.level = 1
p = tf.add_paragraph()
p.text = "Mijozlarga buyurtma holati o'zgarganda avtomatik bot orqali xabar yuboriladi"
p.level = 1
p = tf.add_paragraph()
p.text = "Xavfsizlik: Ro'yxatdan o'tmagan foydalanuvchilar ma'lumotlarni ko'rolmaydi"
p.level = 1

add_logo(slide, logo_2, Inches(0.5), Inches(5.5), Inches(1.5))

# --- Slide 5: Tizim qanday ishlaydi ---
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Tizim arxitekturasi va jarayonlar"
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = "1. Menejer mijozni va buyurtmani CRM ga kiritadi"
p.level = 0
p = tf.add_paragraph()
p.text = "2. Tizim avtomatik ravishda PDF formatida rasmiy shartnoma generatsiya qiladi"
p.level = 0
p = tf.add_paragraph()
p.text = "3. Ishlab chiqarish va Ombor bo'limi buyurtma bo'yicha ishlarni boshlaydi"
p.level = 0
p = tf.add_paragraph()
p.text = "4. Mijoz o'z portali yoki Telegram bot orqali barcha qadamlarni ko'rib boradi"
p.level = 0
p = tf.add_paragraph()
p.text = "5. Rahbariyat Dashboard (Analitika) orqali kirm-chiqimlarni va foydani nazorat qiladi"
p.level = 0

# Save presentation
prs.save(r'C:\Users\rahmo\OneDrive\Desktop\e-mebel\E_Mebel_Taqdimot.pptx')
print("Presentation successfully created at C:\\Users\\rahmo\\OneDrive\\Desktop\\e-mebel\\E_Mebel_Taqdimot.pptx")
